
import re
import uuid
from pptx import Presentation
from collections import Counter
from .database import SessionLocal
from .models import LectureSession
sessions: dict[str, dict] = {}

FOOTER_PATTERNS = [
    r"https?://\S+",
    r"\b\S+@\S+\b",
    r"©|copyright",
    r"\b(all rights reserved)\b",
]
FOOTER_RE = re.compile("|".join(FOOTER_PATTERNS), re.I)

def _clean_lines(lines: list[str]) -> list[str]:
    cleaned = []
    for line in lines:
        if not line:
            continue
        line = re.sub(r"[\u200B-\u200D\uFEFF\x00-\x1F\x7F]", " ", line)
        line = re.sub(r"\s+", " ", line).strip()
        if not line:
            continue
        
        if re.fullmatch(r"\d{1,3}", line) or re.match(r"^\s*slide\s+\d+\s*$", line, re.I):
            continue
        if FOOTER_RE.search(line):
            continue
        cleaned.append(line)
    return cleaned

def extract_text_by_slide(file):
    prs = Presentation(file)
    slides = []

    raw_per_slide = []
    for slide in prs.slides:
        lines = []
        for shape in slide.shapes:
            
            if hasattr(shape, "text_frame") and shape.text_frame:
                for p in shape.text_frame.paragraphs:
                    if p.text and p.text.strip():
                        lines.append(p.text.strip())
            
            if getattr(shape, "shape_type", None) == 19: 
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text and cell.text.strip():
                            lines.append(cell.text.strip())
            
            if hasattr(shape, "shapes"):
                for s in shape.shapes:
                    if hasattr(s, "text_frame") and s.text_frame:
                        for p in s.text_frame.paragraphs:
                            if p.text and p.text.strip():
                                lines.append(p.text.strip())
        raw_per_slide.append(lines)

    
    all_lines = [ln for lines in raw_per_slide for ln in lines]
    norm = lambda t: re.sub(r"\s+", " ", t.strip().lower())
    counts = Counter(norm(t) for t in all_lines if t and len(t) > 10)
    common = {t for t, c in counts.items() if c >= 3}  # appears on ≥3 slides

    for i, lines in enumerate(raw_per_slide, start=1):
        title = ""
        
        lines = _clean_lines(lines)
        lines = [ln for ln in lines if norm(ln) not in common]
        if lines:
            title = lines[0]
        body = "\n".join(lines[1:]).strip() if len(lines) > 1 else ""
        slides.append({"page": i, "title": title or f"Slide {i}", "text": body})
    return slides

def create_session(
    pptx_text: str,
    summary_text: str,
    slides_payload: list[dict] | None = None,
    user_id: int | None = None,
) -> str:
    
    sid = str(uuid.uuid4())

    db = SessionLocal()
    try:
        db_sess = LectureSession(
            id=sid,
            user_id=user_id,
            pptx_text=pptx_text,
            summary_text=summary_text,
            slides_payload=slides_payload or [],
        )
        db.add(db_sess)
        db.commit()
    finally:
        db.close()

    
    sessions[sid] = {
        "user_id": user_id,
        "pptx_text": pptx_text,
        "summary": summary_text,
        "slides": slides_payload or [],
        "chat_history": [],
    }
    return sid

def get_session(session_id: str) -> dict | None:
    
    if session_id in sessions:
        return sessions[session_id]

    db = SessionLocal()
    try:
        db_sess = db.query(LectureSession).filter(LectureSession.id == session_id).first()
        if not db_sess:
            return None
        
        sess = {
            "user_id": db_sess.user_id,
            "pptx_text": db_sess.pptx_text,
            "summary": db_sess.summary_text or "",
            "slides": db_sess.slides_payload or [],
            "chat_history": [],
        }
        sessions[session_id] = sess
        return sess
    finally:
        db.close()
